import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Fix for compatibility with older python-pptx versions and python 3.10+
import sys
if not hasattr(collections, 'Container'):
    collections.Container = collections.abc.Container
if not hasattr(collections, 'Sized'):
    collections.Sized = collections.abc.Sized
if not hasattr(collections, 'Iterable'):
    collections.Iterable = collections.abc.Iterable
if not hasattr(collections, 'Sequence'):
    collections.Sequence = collections.abc.Sequence
if not hasattr(collections, 'Mapping'):
    collections.Mapping = collections.abc.Mapping

def create_presentation():
    prs = Presentation()

    # Define common styles
    def set_title(slide, text):
        title = slide.shapes.title
        title.text = text
        for p in title.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

    def add_bullet_points(slide, points):
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for pt in points:
            p = tf.add_paragraph()
            p.text = pt
            p.level = 0
            p.space_after = Pt(14)
            
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Docu-Sync"
    subtitle.text = "The Next-Gen Collaborative Editor\nReal-Time Sync, Version Control & Time-Travel"

    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "The Problem Statement")
    add_bullet_points(slide, [
        "Need for highly concurrent, version-controlled real-time editing.",
        "Traditional editors fail at maintaining granular history.",
        "Merging simultaneous edits often leads to conflicts and data loss.",
        "Lack of 'time-travel' to revert complex changes seamlessly."
    ])

    # Slide 3: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "Our Solution: Docu-Sync")
    add_bullet_points(slide, [
        "A Google Docs-style rich text editor built for the modern web.",
        "Features a 'Git-like' versioning system for document states.",
        "Powered by CRDTs (Conflict-free Replicated Data Types) for zero-conflict synchronization.",
        "Allows real-time visual diff-tracking and user presence."
    ])

    # Slide 4: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "Key Features - Phased Approach")
    add_bullet_points(slide, [
        "Phase 1: Real-Time Foundation (Zero-conflict editing with Yjs & WebRTC)",
        "Phase 2: Visual Diff-Tracking (Real-time user cursors and presence via Yjs Awareness)",
        "Phase 3: Snapshotting & Time-Travel (Save & restore history like Git branches)"
    ])

    # Slide 5: Technical Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "Technical Architecture")
    add_bullet_points(slide, [
        "Frontend: Vite + React + Vanilla CSS (Premium, responsive US/UX)",
        "Core Editor: React-Quill for rich text manipulation",
        "Sync Engine: Yjs CRDT over WebRTC/WebSockets for decentralized sync",
        "Storage: LocalStorage & Firebase for snapshot persistence"
    ])

    # Slide 6: Under the Hood - Data Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "Under the Hood: The Data Flow")
    add_bullet_points(slide, [
        "1. Interception: React-Quill captures keystrokes and generates fractional 'Delta' updates.",
        "2. CRDT Encoding: Yjs (y-quill) mathematically encodes these Deltas into conflict-free binary states.",
        "3. Broadcasting: WebRTC/WebSockets instantly blast this state to all connected peers.",
        "4. Seamless Merging: Peers receive the update, Yjs resolves history conflicts automatically, and updates the local UI.",
        "5. Awareness: The Yjs Awareness protocol pushes transient data (like cursor location) independently."
    ])

    # Slide 7: Time-Travel Magic
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "The 'Time-Travel' Feature")
    add_bullet_points(slide, [
        "Captures complete 'Snapshots' of document states natively using Yjs.",
        "Persists binary updates to Firebase/LocalStorage.",
        "Allows users to instantly jump to previous snapshots without breaking other users' states.",
        "Enables a secure, revertible workspace."
    ])

    # Slide 7: Why CRDTs over OT?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "Technical Decision: Yjs vs OT")
    add_bullet_points(slide, [
        "Yjs (CRDT) decentralizes the conflict resolution (peer-to-peer ready).",
        "More scalable and less server-intensive than Operational Transformation (OT).",
        "Offline-first approach by design.",
        "Handles edge cases like concurrent formatting effortlessly."
    ])

    # Slide 8: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "Thank You!")
    add_bullet_points(slide, [
        "Docu-Sync is ready to revolutionize collaboration.",
        "Questions?"
    ])

    prs.save('Docu_Sync_Presentation.pptx')
    print("Presentation saved as Docu_Sync_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
